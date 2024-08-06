from datetime import datetime
import os, tempfile
from pptx import Presentation
from pptx.util import Cm, Pt
import pptx
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from ZW_PPT_Title_Head import PPT_Title_Head
import MySQL, ZW_PPT_Chart, ZW_PPT_Table, Trans_Pic
import re

def PPTX():
    # 创建一个新的演示文稿
    ppt = Presentation()
    # 设置幻灯片的宽度和高度  
    ppt.slide_width = Cm(22.15)
    ppt.slide_height = Cm(77)
    # 设置幻灯片对象模型索引6，空白模板
    slide_layout = ppt.slide_layouts[6]
    # 添加一个空白幻灯片
    slide = ppt.slides.add_slide(slide_layout)
    # 添加抬头图片
    image = slide.shapes.add_picture('./Pic/图片1.jpg', Cm(0), Cm(0), Cm(22.15), Cm(4.2))
    # 添加结尾图片
    image = slide.shapes.add_picture('./Pic/图片2.png', Cm(0), Cm(72), Cm(22.15), Cm(4))
    # 添加标题文本框
    PPT_Title_Head()
    image = slide.shapes.add_picture('./Pic/Title_Head.png', Cm(0), Cm(4.3), Cm(22.15), Cm(3.5))
    # 添加装维工单情况图片
    image = slide.shapes.add_picture('./Pic/装维工单情况.png', Cm(0.385), Cm(8), Cm(21.38), Cm(1.6))

    ZW_PPT_Chart.Draw_Chart_1()
    ZW_PPT_Chart.Draw_Chart_2()

    ZW_PPT_Table.Draw_Table_1()
    ZW_PPT_Table.Draw_Table_2()


    # 添加Chart_1图片
    image = slide.shapes.add_picture('./Pic/Chart_1.png', Cm(0), Cm(32.9), Cm(10), Cm(5.5))
    # 添加Chart_2图片
    image = slide.shapes.add_picture('./Pic/Chart_2.png', Cm(12.15), Cm(32.9), Cm(10), Cm(5.5))

    # 添加Table_1图片
    image = slide.shapes.add_picture('./Pic/Table_1.png', Cm(0.5), Cm(10.85), Cm(21.15), Cm(10.15))
    # 添加Table_2图片
    image = slide.shapes.add_picture('./Pic/Table_2.png', Cm(0.5), Cm(39.25), Cm(21.15), Cm(9.75))
    # 添加Table_3图片
    if MySQL.Table_3.iloc[8, 0] != 0:
        ZW_PPT_Table.Draw_Table_3()
        image = slide.shapes.add_picture('./Pic/Table_3.png', Cm(0.5), Cm(49.85), Cm(21.15), Cm(10.05))
    else:
        # 创建文本框并添加文本
        textbox = slide.shapes.add_textbox(left=Cm(0.5), top=Cm(50), width=Cm(21.15), height=Cm(10.05))
        tf = textbox.text_frame
        # 添加文本，并设置格式
        p = tf.paragraphs[0]  # 添加段落
        run = p.add_run()  # 创建一个Run对象，不立即添加文本
        run.text = '当月不满意、质差工单尚未更新，暂时不提供数据'  # 设置Run对象的文本
        run.font.bold = True  # 设置文本加粗
        run.font.size = Pt(18)  # 设置字号大小
        run.font.name = '微软雅黑'  # 设置字体
        # 设置文本的对齐方式
        p.alignment = PP_ALIGN.LEFT


    # 第一个文本框
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(10), width=Cm(22.15), height=Cm(0.85))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = '1、装维在途工单统计'  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT
    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(252, 213, 181)  # 设置RGB颜色值，例如蓝色

    # 第二个文本框
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(32), width=Cm(22.15), height=Cm(0.85))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = '2、个人在途工单量较大的装维人员'  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT
    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(252, 213, 181)  # 设置RGB颜色值，例如蓝色

    # 第三个文本框
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(38.4), width=Cm(22.15), height=Cm(0.85))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = '3、 全市在途工单超时趋势表'  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT
    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(252, 213, 181)  # 设置RGB颜色值，例如蓝色

    # 第四个文本框
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(49), width=Cm(22.15), height=Cm(0.85))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = '4、 全市质差工单整治成功率'  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT
    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(252, 213, 181)  # 设置RGB颜色值，例如蓝色

    # 第五个文本框
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(61.3), width=Cm(22.15), height=Cm(0.85))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = '5、 全市参评“百千工程”分支局“装移机履约及时率”未达标单位'  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT
    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(252, 213, 181)  # 设置RGB颜色值，例如蓝色

    # 文本Txt_2
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0.575), top=Cm(21), width=Cm(21), height=Cm(10.5))
    tf = textbox.text_frame
    # 添加文本，并设置格式


    text = MySQL.Txt_2
    pattern = r'\n'
    results = re.split(pattern, text)



    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[0] + "\n" + results[1]  # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体

    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[2] + "\n" + results[3] + "\n" + results[4] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    run.font.color.rgb = RGBColor(255, 0, 0)  

    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[5] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体


    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[6] + "\n" + results[7] + "\n" + results[8] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    run.font.color.rgb = RGBColor(255, 0, 0) 

    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[9] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体

    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[10] + "\n" + results[11] + "\n" + results[12] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    run.font.color.rgb = RGBColor(255, 0, 0) 

    p = tf.add_paragraph()
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = results[13] # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(16)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体

    line = textbox.line
    line.fill.solid()  # 确保线条填充为纯色
    line.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 红色
    line.width = Pt(1)  # 设置线条宽度
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT


    # 文本Txt_6
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(59.9), width=Cm(22.15), height=Cm(1))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = MySQL.Txt_6  # 设置Run对象的文本
    run.font.bold = False  # 设置文本加粗
    run.font.size = Pt(14)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体

    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT

    # 设置文本框的填充颜色（实际上是设置形状的填充颜色）
    fill = textbox.fill
    fill.solid()  # 设置填充为纯色
    fill.fore_color.rgb = RGBColor(253, 234, 218)  # 设置RGB颜色值，例如蓝色

    # 文本Txt_7
    # 创建文本框并添加文本
    textbox = slide.shapes.add_textbox(left=Cm(1.2), top=Cm(62.75), width=Cm(19.75), height=Cm(9))
    tf = textbox.text_frame
    # 添加文本，并设置格式
    p = tf.paragraphs[0]  # 添加段落
    run = p.add_run()  # 创建一个Run对象，不立即添加文本
    run.text = MySQL.Txt_7  # 设置Run对象的文本
    run.font.bold = True  # 设置文本加粗
    run.font.size = Pt(24)  # 设置字号大小
    run.font.name = '微软雅黑'  # 设置字体
    line = textbox.line
    line.fill.solid()  # 确保线条填充为纯色
    line.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 红色
    line.width = Pt(1)  # 设置线条宽度
    # 设置文本的对齐方式
    p.alignment = PP_ALIGN.LEFT

    textbox = slide.shapes.add_textbox(left=Cm(0), top=Cm(72.8), width=Cm(25), height=Cm(1))
    tf = textbox.text_frame



    ppt.save('./Output/ZW_Daily_Report.pptx')
    Trans_Pic.Trans()

PPTX()